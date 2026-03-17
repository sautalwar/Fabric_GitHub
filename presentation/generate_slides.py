from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PRIMARY_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x20, 0x50)
ACCENT_CYAN = RGBColor(0x50, 0xE6, 0xFF)
SUCCESS_GREEN = RGBColor(0x10, 0x7C, 0x10)
LIGHT_GREEN = RGBColor(0xE7, 0xF4, 0xE4)
WARNING_RED = RGBColor(0xD1, 0x34, 0x38)
LIGHT_RED = RGBColor(0xFE, 0xF0, 0xF1)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xE8, 0xEA, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)
FONT_NAME = "Aptos"

OUTPUT_PATH = Path(r"C:\Users\sautalwar\Downloads\repos\Fabric_GitHub\presentation\Fabric_CICD_Demo.pptx")
SCREENSHOTS = Path(r"C:\Users\sautalwar\Downloads\repos\Fabric_GitHub\presentation\screenshots")


def set_slide_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color



def add_accent_line(slide):
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(1.08), Inches(12.3), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_BLUE
    line.line.fill.background()



def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(11.6), Inches(0.6))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_NAME
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.18), Inches(11.2), Inches(0.45))
        stf = subtitle_box.text_frame
        sp = stf.paragraphs[0]
        srun = sp.add_run()
        srun.text = subtitle
        srun.font.name = FONT_NAME
        srun.font.size = Pt(20)
        srun.font.color.rgb = PRIMARY_BLUE

    add_accent_line(slide)



def set_shape_style(shape, fill_color=None, line_color=None, line_width=1.5):
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)



def add_text_in_shape(shape, lines, font_size=18, color=BLACK, bold_first=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)

    for idx, text in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold_first and idx == 0
        if idx > 0:
            p.space_before = Pt(2)



def add_bullets(slide, items, left, top, width, height, font_size=20, color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.name = FONT_NAME
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box



def add_footer(slide, text, color=GRAY):
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(11)
    run.font.color.rgb = color



def add_labeled_panel(slide, left, top, width, height, title, title_fill, border_color, items, symbol):
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.55))
    set_shape_style(header, title_fill, title_fill)
    add_text_in_shape(header, [title], font_size=22, color=WHITE, bold_first=True)

    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top + Inches(0.62), width, height - Inches(0.62))
    set_shape_style(body, WHITE, border_color)

    text_items = [f"{symbol} {item}" for item in items]
    add_bullets(slide, text_items, left + Inches(0.15), top + Inches(0.78), width - Inches(0.3), height - Inches(0.9), font_size=19, color=DARK_BLUE)



def add_arrow(slide, x1, y1, x2, y2, color, width=2.5):
    arrow_head = Inches(0.18)
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2 - arrow_head // 2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)

    head = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        x2 - arrow_head,
        y2 - arrow_head // 2,
        arrow_head,
        arrow_head,
    )
    head.rotation = 90
    set_shape_style(head, color, color)
    return connector



def add_slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Microsoft Fabric — SDLC & CI/CD Overview", "Automated deployment across Dev → UAT → Prod")

    banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.55), Inches(3.9), Inches(0.65))
    set_shape_style(banner, ACCENT_CYAN, ACCENT_CYAN)
    add_text_in_shape(banner, ["Digital Realty Demo"], font_size=22, color=DARK_BLUE, bold_first=True)

    bullets = [
        "Fabric Workspaces as environments",
        "Git Integration for version control",
        "Deployment Pipelines for promotion",
        "REST APIs for automation",
    ]
    add_bullets(slide, bullets, Inches(0.9), Inches(1.9), Inches(6.2), Inches(3.4), font_size=22, color=BLACK)

    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.35), Inches(4.2), Inches(2.75))
    set_shape_style(panel, RGBColor(0xF2, 0xF8, 0xFC), PRIMARY_BLUE)
    add_text_in_shape(
        panel,
        [
            "Core Flow",
            "Dev workspace builds assets",
            "Git captures supported items",
            "Deployment pipeline promotes changes",
            "APIs orchestrate end-to-end automation",
        ],
        font_size=19,
        color=DARK_BLUE,
        bold_first=True,
        align=PP_ALIGN.LEFT,
    )

    add_footer(slide, "Digital Realty — Fabric CI/CD Workshop")



def add_slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "What Git Integration Tracks")

    tracked = [
        "Notebooks",
        "Data Pipelines",
        "Semantic Models",
        "Reports/Dashboards",
        "Spark Job Definitions",
        "ML Models",
    ]
    not_tracked = [
        "Lakehouse Tables",
        "Table Schemas & DDL",
        "Views & Stored Procedures",
        "Data Files",
        "SQL Endpoint Objects",
    ]

    add_labeled_panel(slide, Inches(0.8), Inches(1.7), Inches(5.7), Inches(4.6), "✅ Tracked", PRIMARY_BLUE, PRIMARY_BLUE, tracked, "✅")
    add_labeled_panel(slide, Inches(6.9), Inches(1.7), Inches(5.7), Inches(4.6), "❌ Not Tracked", WARNING_RED, WARNING_RED, not_tracked, "❌")

    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.45))
    set_shape_style(note, RGBColor(0xF7, 0xFB, 0xFE), LIGHT_GRAY)
    add_text_in_shape(note, ["Source: Microsoft Learn — Lakehouse Git Integration"], font_size=14, color=GRAY)



def add_slide_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Schema Gap")

    dev = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(2.35), Inches(2.7), Inches(1.45))
    set_shape_style(dev, RGBColor(0xF2, 0xF8, 0xFC), PRIMARY_BLUE)
    add_text_in_shape(dev, ["Dev Lakehouse", "Tables + schema", "Created in workspace"], font_size=20, color=DARK_BLUE, bold_first=True)

    git_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.05), Inches(2.35), Inches(2.7), Inches(1.45))
    set_shape_style(git_box, RGBColor(0xF5, 0xF5, 0xF5), GRAY)
    add_text_in_shape(git_box, ["Git", "Artifacts only", "No table metadata"], font_size=20, color=GRAY, bold_first=True)

    uat = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(2.35), Inches(2.8), Inches(1.45))
    set_shape_style(uat, RGBColor(0xFB, 0xFB, 0xFB), PRIMARY_BLUE)
    add_text_in_shape(uat, ["UAT Lakehouse", "Empty / missing tables", "Schema must be rebuilt"], font_size=20, color=DARK_BLUE, bold_first=True)

    add_arrow(slide, Inches(3.65), Inches(3.08), Inches(5.0), Inches(3.08), PRIMARY_BLUE)
    add_arrow(slide, Inches(7.75), Inches(3.08), Inches(9.05), Inches(3.08), WARNING_RED)

    gap = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(7.55), Inches(1.85), Inches(1.15), Inches(0.75))
    set_shape_style(gap, LIGHT_RED, WARNING_RED, line_width=2)
    add_text_in_shape(gap, ["Gap"], font_size=20, color=WARNING_RED, bold_first=True)

    key_point = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(4.5), Inches(10.95), Inches(1.15))
    set_shape_style(key_point, LIGHT_RED, WARNING_RED)
    add_text_in_shape(
        key_point,
        ["Schema changes made in Dev Lakehouse do NOT propagate through Git or Deployment Pipelines"],
        font_size=24,
        color=WARNING_RED,
        bold_first=True,
    )

    caption = slide.shapes.add_textbox(Inches(1.15), Inches(5.9), Inches(10.7), Inches(0.5))
    ctf = caption.text_frame
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = "Result: promoted environments receive notebooks and pipelines, but not the underlying lakehouse schema."
    cr.font.name = FONT_NAME
    cr.font.size = Pt(18)
    cr.font.color.rgb = GRAY



def add_slide_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "The Solution — Schema-as-Code")

    git_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.75), Inches(2.45), Inches(1.1))
    set_shape_style(git_box, RGBColor(0xF5, 0xF5, 0xF5), GRAY)
    add_text_in_shape(git_box, ["Migration Scripts in Git", "SQL + PySpark files"], font_size=19, color=GRAY, bold_first=True)

    pipeline = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(1.75), Inches(2.45), Inches(1.1))
    set_shape_style(pipeline, LIGHT_GREEN, SUCCESS_GREEN)
    add_text_in_shape(pipeline, ["CI/CD Pipeline", "GitHub Actions / Azure DevOps"], font_size=19, color=SUCCESS_GREEN, bold_first=True)

    runner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.45), Inches(1.75), Inches(2.55), Inches(1.1))
    set_shape_style(runner, RGBColor(0xF2, 0xF8, 0xFC), PRIMARY_BLUE)
    add_text_in_shape(runner, ["Migration Runner Notebook", "Idempotent execution engine"], font_size=19, color=DARK_BLUE, bold_first=True)

    envs = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.4), Inches(1.55), Inches(2.9), Inches(1.5))
    set_shape_style(envs, RGBColor(0xF3, 0xFB, 0xF3), SUCCESS_GREEN)
    add_text_in_shape(envs, ["Dev / UAT / Prod", "Lakehouse schema applied", "History recorded"], font_size=19, color=SUCCESS_GREEN, bold_first=True)

    add_arrow(slide, Inches(3.2), Inches(2.3), Inches(3.55), Inches(2.3), SUCCESS_GREEN)
    add_arrow(slide, Inches(6.05), Inches(2.3), Inches(6.4), Inches(2.3), SUCCESS_GREEN)
    add_arrow(slide, Inches(9.0), Inches(2.3), Inches(9.35), Inches(2.3), SUCCESS_GREEN)

    bullet_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(3.45), Inches(11.35), Inches(2.65))
    set_shape_style(bullet_box, WHITE, LIGHT_GRAY)
    bullets = [
        "Version-controlled SQL & PySpark migration scripts",
        "Idempotent migration runner notebook",
        "Automated CI/CD pipeline (Azure DevOps or GitHub Actions)",
        "Audit trail via _migration_history table",
    ]
    add_bullets(slide, bullets, Inches(1.2), Inches(3.72), Inches(10.85), Inches(2.15), font_size=20, color=BLACK)



def add_slide_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Demo Roadmap")

    roadmap_items = [
        "1. Environment Setup (3 workspaces + Deployment Pipeline)",
        "2. Standard CI/CD Flow (notebooks through Git)",
        "3. The Problem (tables don't propagate)",
        "4. Schema-as-Code Solution (migration scripts + runner)",
        "5. CI/CD Pipeline Automation (Azure DevOps + GitHub Actions)",
        "6. REST API Deep Dive (live API calls)",
    ]
    add_bullets(slide, roadmap_items, Inches(1.0), Inches(1.65), Inches(10.8), Inches(4.55), font_size=21, color=DARK_BLUE)

    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(6.05), Inches(4.0), Inches(0.75))
    set_shape_style(callout, ACCENT_CYAN, PRIMARY_BLUE)
    add_text_in_shape(callout, ["Let's see it in action! →"], font_size=22, color=DARK_BLUE, bold_first=True)



def add_slide_environments(prs):
    """Slide showing real Fabric environment screenshots."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Environment Setup — 3 Workspaces", "Dev → UAT → Prod with DataCenterLakehouse")

    img_specs = [
        ("dev_lakehouse.png", "Dev Workspace", Inches(0.4), Inches(1.65), Inches(4.1)),
        ("uat_lakehouse.png", "UAT Workspace", Inches(4.6), Inches(1.65), Inches(4.1)),
        ("prod_lakehouse.png", "Prod Workspace", Inches(8.8), Inches(1.65), Inches(4.1)),
    ]
    for filename, label, left, top, width in img_specs:
        img_path = SCREENSHOTS / filename
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), left, top, width)
            label_box = slide.shapes.add_textbox(left, top + Inches(3.35), width, Inches(0.4))
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.name = FONT_NAME
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY_BLUE

    add_footer(slide, "Digital Realty — Fabric CI/CD Workshop")


def add_slide_pipeline(prs):
    """Slide showing the completed deployment pipeline screenshot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Deployment Pipeline — Dev → UAT → Prod", "DigitalRealty-Pipeline with all stages configured")

    img_path = SCREENSHOTS / "deployment_pipeline_complete.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.8), Inches(1.55), Inches(11.7))

    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.5))
    set_shape_style(note, RGBColor(0xF2, 0xF8, 0xFC), PRIMARY_BLUE)
    add_text_in_shape(note, ["Pipeline deploys supported artifacts • Schema requires Schema-as-Code pattern"], font_size=16, color=DARK_BLUE)

    add_footer(slide, "Digital Realty — Fabric CI/CD Workshop")


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Fabric CI/CD Demo"
    prs.core_properties.subject = "Microsoft Fabric SDLC & CI/CD Overview"
    prs.core_properties.author = "GitHub Copilot CLI"
    prs.core_properties.company = "Digital Realty"

    add_slide_one(prs)             # 1. Overview
    add_slide_environments(prs)    # 2. Environment screenshots
    add_slide_two(prs)             # 3. What Git tracks
    add_slide_three(prs)           # 4. The schema gap
    add_slide_four(prs)            # 5. Schema-as-Code solution
    add_slide_pipeline(prs)        # 6. Deployment pipeline screenshot
    add_slide_five(prs)            # 7. Demo roadmap

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH



if __name__ == "__main__":
    output = build_presentation()
    print(f"Created presentation: {output}")
